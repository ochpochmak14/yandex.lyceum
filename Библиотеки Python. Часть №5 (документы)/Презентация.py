import random
from pptx import Presentation
from pptx.util import Pt


def get_method_description(method_name):
    docstring = getattr(random, method_name).__doc__
    return docstring if docstring else "Описание отсутствует."


methods = ["randint", "random", "choice", "shuffle", "uniform"]

presentation = Presentation()
font_name = "Courier New"
font_size = Pt(18)

for method in methods:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = f"Метод random.{method}"
    content.text = get_method_description(method)

    for paragraph in content.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = font_size

presentation.save("res.pptx")
